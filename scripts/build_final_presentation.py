"""Build and validate the editable SynthDet final PowerPoint presentation."""

from __future__ import annotations

import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"
OUTPUT = OUT_DIR / "SynthDet_Final_Presentation.pptx"
FIGURES = ROOT / "reports" / "final" / "figures"
LOGO = ROOT / "apps" / "web" / "public" / "brand" / "synthdet-logo.png"

NAVY = RGBColor(8, 18, 36)
PANEL = RGBColor(18, 32, 57)
BLUE = RGBColor(54, 127, 245)
CYAN = RGBColor(29, 190, 205)
WHITE = RGBColor(246, 249, 255)
MUTED = RGBColor(164, 182, 209)
GREEN = RGBColor(31, 190, 139)
AMBER = RGBColor(245, 184, 48)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def rtl(paragraph) -> None:
    paragraph.alignment = PP_ALIGN.RIGHT
    paragraph._p.get_or_add_pPr().set("rtl", "1")


def set_background(slide, color: RGBColor = NAVY) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 22,
    color: RGBColor = WHITE,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.RIGHT,
):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    if align == PP_ALIGN.RIGHT:
        rtl(paragraph)
    run = paragraph.runs[0]
    run.font.name = "Tajawal"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def title(slide, text: str, number: int) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    textbox(slide, text, 5.15, 0.42, 7.55, 0.7, size=29, bold=True)
    textbox(slide, f"{number:02d}", 0.5, 0.48, 0.6, 0.4, size=13, color=MUTED, align=PP_ALIGN.LEFT)


def footer(slide) -> None:
    textbox(
        slide,
        "sprint5-final-20260720-v1 · نتائج مختومة",
        0.5,
        7.05,
        4.2,
        0.28,
        size=9,
        color=MUTED,
        align=PP_ALIGN.LEFT,
    )


def bullets(
    slide,
    items: list[str],
    *,
    top: float = 1.55,
    left: float = 5.8,
    width: float = 6.55,
    height: float = 4.8,
    size: int = 21,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {item}"
        rtl(paragraph)
        paragraph.space_after = Pt(13)
        run = paragraph.runs[0]
        run.font.name = "Tajawal"
        run.font.size = Pt(size)
        run.font.color.rgb = WHITE


def add_image(
    slide, path: Path, left: float, top: float, width: float, height: float | None = None
) -> None:
    if not path.is_file():
        raise FileNotFoundError(path)
    kwargs = {"width": Inches(width)}
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def panel(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    value: str,
    label: str,
    color: RGBColor = BLUE,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = color
    textbox(
        slide, value, left + 0.15, top + 0.25, width - 0.3, 0.62, size=27, bold=True, color=color
    )
    textbox(slide, label, left + 0.15, top + 0.95, width - 0.3, 0.5, size=13, color=MUTED)


def add_content_slide(
    prs: Presentation, number: int, heading: str, items: list[str], image: str | None = None
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    title(slide, heading, number)
    if image:
        add_image(slide, FIGURES / image, 0.55, 1.35, 5.55, 4.95)
        bullets(slide, items, left=6.35, width=6.35, size=18)
    else:
        bullets(slide, items, left=1.0, width=11.3, size=22)
    footer(slide)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_image(slide, LOGO, 0.65, 0.55, 2.0, 2.0)
    textbox(slide, "سينث دِت — SynthDet", 4.1, 1.2, 8.2, 0.8, size=38, bold=True)
    textbox(
        slide,
        "تعزيز البيانات الاصطناعية لكشف الأجسام باستخدام YOLO",
        3.0,
        2.15,
        9.3,
        0.62,
        size=25,
        color=CYAN,
        bold=True,
    )
    textbox(
        slide,
        "قياس فجوة المجال بين البيانات الاصطناعية والحقيقية",
        3.0,
        2.95,
        9.3,
        0.55,
        size=22,
        color=MUTED,
    )
    panel(slide, 3.1, 4.25, 2.7, 1.45, "5", "أنظمة تدريب مكتملة", GREEN)
    panel(slide, 5.95, 4.25, 2.7, 1.45, "68", "صورة اختبار محمية", AMBER)
    panel(slide, 8.8, 4.25, 2.7, 1.45, "0.2119", "قيمة المقياس الأساسي", BLUE)
    footer(slide)

    add_content_slide(
        prs,
        2,
        "المشكلة والدافع",
        [
            "تعليم الصور الحقيقية مكلف ومحدود لبعض الفئات.",
            "الصور المركبة قد تزيد التنوع لكنها تحمل فجوة مظهرية.",
            "المطلوب قياس النقل إلى صور حقيقية جديدة تحت مقارنة عادلة.",
        ],
    )
    add_content_slide(
        prs,
        3,
        "الهدف وأسئلة البحث",
        [
            "تثبيت المعمارية والميزانية والتحقق والإعدادات.",
            "تغيير نسبة الحقيقي/الاصطناعي فقط: 0، 25، 50، 75، 100%.",
            "ترتيب النماذج على اختبار حقيقي محمي وفق قاعدة معلنة مسبقًا.",
        ],
    )
    add_content_slide(
        prs,
        4,
        "سير العمل العلمي",
        [
            "كل مرحلة تربط هويتها وبصمتها بالمرحلة التالية.",
            "لا فتح للاختبار قبل تجميد العقد ودفعه.",
            "النتائج والمنتج مشتقان من المخرجات المختومة.",
        ],
        "figure_05_workflow.png",
    )
    add_content_slide(
        prs,
        5,
        "البيانات وسبع فئات",
        [
            "635 صورة مقبولة و4,784 جسمًا بعد تحقق صارم.",
            "fish، jellyfish، penguin، puffin، shark، starfish، stingray.",
            "عدم توازن فئوي ومشاهد انعكاس وازدحام وإضاءة صناعية.",
        ],
    )
    add_content_slide(
        prs,
        6,
        "التقسيم الثاني ومنع التسرب",
        [
            "427 تدريب / 140 تحقق / 68 اختبار محمي.",
            "تقسيم مجموعات مصدر غير قابلة للتجزئة ببذرة 42.",
            "تصحيح فئة البطريق: أربع صور اختبار فقط؛ عدم يقين مرتفع.",
        ],
    )
    add_content_slide(
        prs,
        7,
        "التوليد الاصطناعي",
        [
            "427 مركب نسخ ولصق من التدريب فقط، وليست مشاهد مرسومة بالكامل.",
            "798 جسمًا ملصقًا مع أقنعة وفلاتر جودة وتحويلات محدودة.",
            "بصمة وتتبّع لكل مصدر وخلفية وتحويل ومحاولة وضع.",
        ],
    )
    add_content_slide(
        prs,
        8,
        "تصميم التجارب",
        [
            "427 صورة في كل نظام؛ تحقق حقيقي مشترك.",
            "اقتران تكميلي: القماش الحقيقي أو مشتقه الاصطناعي مرة واحدة.",
            "المتغير العلمي الوحيد هو تركيب بيانات التدريب.",
        ],
    )
    add_content_slide(
        prs,
        9,
        "بروتوكول التدريب",
        [
            "نموذج الكشف الصغير، حجم 640، خمسون حقبة، ودفعة 16.",
            "بطاقة تسريع موحدة وبذرة 42 وإعدادات تحسين ثابتة.",
            "اكتملت خمسة best/last؛ وصول الاختبار أثناء التدريب = صفر.",
        ],
    )
    add_content_slide(
        prs,
        10,
        "نتائج التحقق — غير نهائية",
        [
            "قيم المقياس بالترتيب: 0.3011، 0.3263، 0.3193، 0.3308، 0.3125.",
            "نظام 75% حقيقي تصدر التحقق، لكنه لم يتصدر الاختبار.",
            "لم يُختر الفائز قبل اكتمال النماذج الخمسة.",
        ],
        "figure_02_domain_gap.png",
    )
    add_content_slide(
        prs,
        11,
        "نتائج الاختبار النهائي",
        [
            "الحقيقي فقط في المرتبة الأولى؛ قيمة المقياس 0.211920.",
            "خليط 50% حقيقي هو الأقوى بين الخلطات؛ قيمته 0.198121.",
            "الترتيب الكامل مُحدد بقاعدة ما قبل الاختبار.",
        ],
        "figure_01_final_metrics.png",
    )
    add_content_slide(
        prs,
        12,
        "حسب الفئة والحجم",
        [
            "متصدر الفئات يختلف؛ لا نظام يهيمن على الجميع.",
            "الحقيقي فقط أفضل للصغير والكبير؛ خليط 50% للمتوسط.",
            "الأجسام الصغيرة ضعيفة والبطريق محدود بأربع صور.",
        ],
        "figure_03_per_class_heatmap.png",
    )
    add_content_slide(
        prs,
        13,
        "فجوة المجال",
        [
            "الحقيقي فقط أعلى من الاصطناعي فقط بمقدار 0.042982.",
            "الفرق النسبي 25.44%، وكل الخلطات تتجاوز الاصطناعي فقط.",
            "المنحنى غير رتيب؛ لا ادعاء نسبة مثلى أو أثر سببي عام.",
        ],
        "figure_02_domain_gap.png",
    )
    add_content_slide(
        prs,
        14,
        "تحليل الأخطاء",
        [
            "235 حالة مختارة بقواعد حتمية بعد ختم الحملة.",
            "إيجابيات وسلبيات كاذبة وتموضع والتباس واختلافات ونجاحات ممثلة.",
            "بكسلات الاختبار لا تُنشر؛ المستودع يتتبع البيانات الوصفية فقط.",
        ],
    )
    add_content_slide(
        prs,
        15,
        "اللوحة العربية وخدمة FastAPI",
        [
            "Next.js RTL وTajawal ووضعا ضوء/ظلام واستجابة كاملة.",
            "ثلاثة أوضاع بيانات صريحة بلا رجوع صامت إلى العرض التجريبي.",
            "استدلال كسول محمي بالبصمة وصورة معلّمة ونتيجة منظمة.",
        ],
    )
    add_content_slide(
        prs,
        16,
        "قابلية إعادة الإنتاج",
        [
            "هويات ثابتة للتقسيم والتوليد والتجارب والتدريب والعقد والحملة.",
            "خمس بصمات نتائج وعشر بصمات تنبؤات متحققة.",
            "أوامر لويندوز ولينكس وحزمة محلية دون نشر الأوزان أو الاختبار.",
        ],
    )
    add_content_slide(
        prs,
        17,
        "القيود",
        [
            "68 صورة اختبار، أربع صور بطريق، وعدم توازن فئوي.",
            "مجموعة ومعمارية وبذرة ومولد نسخ ولصق واحد.",
            "لا فواصل ثقة؛ النتائج ارتباطية ومحدودة الصلاحية الخارجية.",
        ],
    )
    add_content_slide(
        prs,
        18,
        "الخلاصة والعمل المستقبلي",
        [
            "التوصية المجمدة: الحقيقي فقط؛ أقوى خليط: 50% حقيقي.",
            "الاصطناعي ينقل معرفة لكنه لا يغلق الفجوة في هذا التصميم.",
            "مستقبلًا: بذور ومجموعات ومعماريات ومولدات أقوى وعدم يقين.",
        ],
    )
    return prs


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation = build()
    presentation.save(OUTPUT)
    reopened = Presentation(OUTPUT)
    if len(reopened.slides) != 18:
        raise RuntimeError(f"Expected 18 slides, found {len(reopened.slides)}")
    empty = [index + 1 for index, slide in enumerate(reopened.slides) if not slide.shapes]
    if empty:
        raise RuntimeError(f"Empty slides: {empty}")
    validation = {
        "schema_version": 1,
        "file": str(OUTPUT.relative_to(ROOT)).replace("\\", "/"),
        "slide_count": len(reopened.slides),
        "sha256": sha256(OUTPUT),
        "official_logo_sha256": sha256(LOGO),
        "protected_test_images_embedded": False,
        "source": "presentation/source/SynthDet_Final_Presentation.md",
        "generator": "scripts/build_final_presentation.py",
        "status": "validated_structure",
    }
    (OUT_DIR / "PRESENTATION_VALIDATION.json").write_text(
        json.dumps(validation, indent=2) + "\n", encoding="utf-8"
    )
    print(json.dumps(validation, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
